# -*- coding: utf-8 -*-
"""Presentation PowerPoint professionnelle AfriMarket — charte graphique du logo."""
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FIG = r"c:\Users\LENOVO\Desktop\PROJET PYTHON\figures\brand"
LOGO = r"c:\Users\LENOVO\Desktop\PROJET PYTHON\assets\logo.png"
DATA = r"c:\Users\LENOVO\Desktop\PROJET PYTHON\data\df_features.csv"
OUT = r"c:\Users\LENOVO\Desktop\PROJET PYTHON\reports\AfriMarket_Presentation_Direction.pptx"

RED = RGBColor(0xC4, 0x2D, 0x1C)
RED_DARK = RGBColor(0x8E, 0x20, 0x13)
GRAY = RGBColor(0x3F, 0x3F, 0x3F)
GRAY_LIGHT = RGBColor(0x8C, 0x8C, 0x8C)
GRAY_PALE = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# --------------------------------------------------------------------------
# DONNEES (recalculees depuis df_features.csv)
# --------------------------------------------------------------------------
df = pd.read_csv(DATA, parse_dates=["date_commande"], encoding="utf-8")
df_valid = df[df["statut_commande"] != "Annulée"].copy()

ca_total = df_valid["chiffre_affaires"].sum()
profit_total = df_valid["profit_net"].sum()
panier_moyen = df_valid["chiffre_affaires"].mean()
taux_annulation = (df["statut_commande"] == "Annulée").mean() * 100
taux_retour = (df["statut_commande"] == "Retournée").mean() * 100
nb_clients = df["id_client"].nunique()

# --------------------------------------------------------------------------
# HELPERS
# --------------------------------------------------------------------------

def set_no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        set_no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=GRAY, bold=False, align=PP_ALIGN.LEFT,
             italic=False, font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=15, color=GRAY, bold_lead=True, space_after=10, marker_color=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    marker_color = marker_color or RED
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = f"●  {item}"
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = color
    return tb


def slide_header(slide, title, kicker=None):
    """Bandeau superieur commun a tous les slides de contenu."""
    add_rect(slide, 0, 0, SW, Inches(0.12), RED)
    slide.shapes.add_picture(LOGO, Inches(0.35), Inches(0.28), height=Inches(0.55))
    if kicker:
        add_text(slide, Inches(1.1), Inches(0.28), Inches(6), Inches(0.3), kicker.upper(),
                  size=11, color=RED, bold=True)
        add_text(slide, Inches(1.1), Inches(0.52), Inches(10.5), Inches(0.55), title,
                  size=24, color=GRAY, bold=True)
    else:
        add_text(slide, Inches(1.1), Inches(0.35), Inches(10.5), Inches(0.6), title,
                  size=24, color=GRAY, bold=True)
    add_rect(slide, Inches(0.35), Inches(1.15), Inches(3.2), Pt(2.2), RED)


def slide_footer(slide, page_num):
    add_text(slide, Inches(0.35), SH - Inches(0.42), Inches(6), Inches(0.3),
              "AfriMarket — Analyse strategique e-commerce", size=9, color=GRAY_LIGHT, italic=True)
    add_text(slide, SW - Inches(1.2), SH - Inches(0.42), Inches(0.8), Inches(0.3),
              str(page_num), size=9, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)


def new_slide(title=None, kicker=None, page_num=None):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    if title:
        slide_header(slide, title, kicker)
    if page_num:
        slide_footer(slide, page_num)
    return slide


def add_picture_fit(slide, path, x, y, max_w, max_h):
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    slide.shapes.add_picture(path, px, py, width=w, height=h)


def add_logo_badge(slide, cx, cy, logo_h, pad=Inches(0.35)):
    """Place le logo (fond blanc opaque) sur un badge blanc arrondi, pour un rendu
    propre sur les fonds sombres (le PNG source n'a pas de transparence)."""
    from PIL import Image
    im = Image.open(LOGO)
    iw, ih = im.size
    lh = logo_h
    lw = int(iw * (lh / ih))
    badge_w, badge_h = lw + pad, lh + pad
    bx, by = int(cx - badge_w / 2), int(cy - badge_h / 2)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, badge_w, badge_h)
    badge.adjustments[0] = 0.12
    badge.fill.solid()
    badge.fill.fore_color.rgb = WHITE
    set_no_line(badge)
    badge.shadow.inherit = False
    lx, ly = int(cx - lw / 2), int(cy - lh / 2)
    slide.shapes.add_picture(LOGO, lx, ly, width=lw, height=lh)


def kpi_card(slide, x, y, w, h, label, value, accent=RED):
    add_rect(slide, x, y, w, h, GRAY_PALE)
    add_rect(slide, x, y, Inches(0.06), h, accent)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.3), Inches(0.5), label,
              size=12, color=GRAY_LIGHT, bold=False)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), w - Inches(0.3), Inches(0.7), value,
              size=24, color=GRAY, bold=True)

# --------------------------------------------------------------------------
# SLIDE 1 — TITRE
# --------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, GRAY)
add_rect(s, 0, SH - Inches(0.35), SW, Inches(0.35), RED)
add_logo_badge(s, SW / 2, Inches(2.35), Inches(2.0))
add_text(s, Inches(1.5), Inches(3.9), Inches(10.33), Inches(0.9),
          "ANALYSE STRATÉGIQUE E-COMMERCE", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.65), Inches(10.33), Inches(0.5),
          "Performance commerciale · Juillet – Décembre 2025", size=16, color=RGBColor(0xD9,0xD9,0xD9),
          align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(1.5), Inches(6.6), Inches(10.33), Inches(0.4),
          "Présenté à la Direction  ·  Data Analyst — AfriMarket", size=12, color=RGBColor(0xCC,0xCC,0xCC),
          align=PP_ALIGN.CENTER)

# --------------------------------------------------------------------------
# SLIDE 2 — AGENDA
# --------------------------------------------------------------------------
s = new_slide("Sommaire", page_num=2)
agenda = [
    "Contexte & objectifs de l'analyse",
    "Méthodologie : audit et nettoyage des données",
    "Indicateurs clés de performance (KPIs)",
    "Analyse par catégorie de produits",
    "Analyse géographique",
    "Analyse marketing (ROI par canal)",
    "Analyse clients & rétention",
    "Synthèse des insights clés",
    "5 recommandations stratégiques",
    "Conclusion & prochaines étapes",
]
col1 = agenda[:5]
col2 = agenda[5:]
for i, item in enumerate(col1):
    add_text(s, Inches(1.1), Inches(1.7 + i*0.75), Inches(0.5), Inches(0.5), f"{i+1}", size=20, color=RED, bold=True)
    add_text(s, Inches(1.6), Inches(1.72 + i*0.75), Inches(5.2), Inches(0.5), item, size=15, color=GRAY)
for i, item in enumerate(col2):
    add_text(s, Inches(7.1), Inches(1.7 + i*0.75), Inches(0.5), Inches(0.5), f"{i+6}", size=20, color=RED, bold=True)
    add_text(s, Inches(7.6), Inches(1.72 + i*0.75), Inches(5.2), Inches(0.5), item, size=15, color=GRAY)

# --------------------------------------------------------------------------
# SLIDE 3 — CONTEXTE
# --------------------------------------------------------------------------
s = new_slide("Contexte & objectifs", kicker="Introduction", page_num=3)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.1),
    "AfriMarket est une entreprise e-commerce panafricaine active dans 8 villes et 4 catégories "
    "(Électronique, Mode, Beauté, Maison). La direction constate des variations de chiffre "
    "d'affaires, un taux de retour préoccupant et des écarts de performance selon les villes.",
    size=16, color=GRAY, line_spacing=1.3)
add_bullets(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(3),
    [
        "Période analysée : 6 mois d'activité commerciale (juillet – décembre 2025)",
        "9 400 commandes exploitables après nettoyage (sur 10 100 lignes brutes)",
        "Objectif : identifier les leviers de rentabilité et fournir des recommandations actionnables",
    ], size=16, space_after=14)

# --------------------------------------------------------------------------
# SLIDE 4 — METHODOLOGIE
# --------------------------------------------------------------------------
s = new_slide("Méthodologie", kicker="Qualité des données", page_num=4)
problems = [
    ("100", "doublons exacts", "supprimés"),
    ("622", "prix invalides (≤0)", "imputés (médiane categorie)"),
    ("600", "remises négatives", "corrigées (valeur absolue)"),
    ("600", "quantités nulles", "supprimées"),
    ("605", "villes mal orthographiées", "corrigées (Kinshassa → Kinshasa)"),
    ("606", "catégories incohérentes", "fusionnées (electronique → Électronique)"),
]
cols = 3
card_w, card_h = Inches(3.7), Inches(1.5)
gap_x, gap_y = Inches(0.25), Inches(0.25)
start_x, start_y = Inches(0.8), Inches(1.55)
for i, (num, prob, action) in enumerate(problems):
    row, col = divmod(i, cols)
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rect(s, x, y, card_w, card_h, GRAY_PALE)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(1.5), Inches(0.6), num, size=28, color=RED, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.75), card_w - Inches(0.4), Inches(0.35), prob, size=13, color=GRAY, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(1.08), card_w - Inches(0.4), Inches(0.35), action, size=11, color=GRAY_LIGHT, italic=True)

# --------------------------------------------------------------------------
# SLIDE 5 — KPIs GLOBAUX
# --------------------------------------------------------------------------
s = new_slide("Performance globale", kicker="Indicateurs clés", page_num=5)
kpis = [
    ("CA TOTAL", f"{ca_total/1e6:,.2f} M$"),
    ("PROFIT NET ESTIMÉ", f"{profit_total/1000:,.0f} k$"),
    ("PANIER MOYEN", f"{panier_moyen:,.0f} $"),
    ("TAUX D'ANNULATION", f"{taux_annulation:.1f} %"),
    ("TAUX DE RETOUR", f"{taux_retour:.1f} %"),
    ("CLIENTS ACTIFS", f"{nb_clients:,}"),
]
card_w, card_h = Inches(3.75), Inches(1.7)
gap = Inches(0.25)
start_x, start_y = Inches(0.8), Inches(1.8)
for i, (label, val) in enumerate(kpis):
    row, col = divmod(i, 3)
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)
    kpi_card(s, x, y, card_w, card_h, label, val)
add_text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8),
    "Le CA mensuel oscille entre 400k$ et 460k$ sans tendance de croissance nette : "
    "l'instabilité constatée par la direction s'explique par les écarts entre catégories et villes (slides suivantes).",
    size=13, color=GRAY_LIGHT, italic=True)

# --------------------------------------------------------------------------
# SLIDE 6 — CATEGORIE (CA + Retour)
# --------------------------------------------------------------------------
s = new_slide("Performance par catégorie", kicker="Quelle catégorie prioriser ?", page_num=6)
add_picture_fit(s, f"{FIG}/01_ca_categorie.png", Inches(0.5), Inches(1.5), Inches(6.1), Inches(4.3))
add_picture_fit(s, f"{FIG}/02_retour_categorie.png", Inches(6.7), Inches(1.5), Inches(6.1), Inches(4.3))
add_rect(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(1.1), GRAY_PALE)
add_text(s, Inches(0.75), Inches(6.15), Inches(11.9), Inches(0.9),
    "Électronique = 74,6 % du CA mais marge la plus faible (15 %) et taux de retour le plus élevé (13,8 %). "
    "Mode et Beauté (marges 45-50 %) sont sous-exploitées : elles ne pèsent que 10,4 % du CA.",
    size=13, color=GRAY, bold=True)

# --------------------------------------------------------------------------
# SLIDE 7 — EVOLUTION MENSUELLE CATEGORIE
# --------------------------------------------------------------------------
s = new_slide("Évolution mensuelle du CA par catégorie", kicker="Tendances", page_num=7)
add_picture_fit(s, f"{FIG}/03_evolution_categorie.png", Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.6))
add_text(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
    "L'Électronique tire l'essentiel de la variabilité mensuelle du CA global — sa stabilisation est la priorité n°1.",
    size=13, color=GRAY_LIGHT, italic=True)

# --------------------------------------------------------------------------
# SLIDE 8 — GEOGRAPHIE (CA + annulation)
# --------------------------------------------------------------------------
s = new_slide("Performance géographique", kicker="Où investir davantage ?", page_num=8)
add_picture_fit(s, f"{FIG}/04_ca_ville.png", Inches(0.5), Inches(1.5), Inches(6.1), Inches(4.5))
add_picture_fit(s, f"{FIG}/05_annulation_ville.png", Inches(6.7), Inches(1.5), Inches(6.1), Inches(4.5))
add_rect(s, Inches(0.5), Inches(6.15), Inches(12.33), Inches(1.0), GRAY_PALE)
add_text(s, Inches(0.75), Inches(6.3), Inches(11.9), Inches(0.8),
    "Douala : plus forte croissance (+76 %) mais taux d'annulation de 12,9 % (vs 0-1 % ailleurs) "
    "— un problème opérationnel local à résoudre avant d'investir davantage.",
    size=13, color=GRAY, bold=True)

# --------------------------------------------------------------------------
# SLIDE 9 — MARKETING ROI
# --------------------------------------------------------------------------
s = new_slide("Performance marketing", kicker="Quel canal mérite plus de budget ?", page_num=9)
add_picture_fit(s, f"{FIG}/06_roi_marketing.png", Inches(1.5), Inches(1.5), Inches(7.5), Inches(4.6))
add_bullets(s, Inches(9.3), Inches(1.9), Inches(3.5), Inches(4),
    [
        "Email : ROI 231x — sous-investi, à scaler",
        "Instagram Ads : meilleure rétention (54%), à maintenir",
        "Influenceur : ROI le plus faible (21,6x) — à réduire",
    ], size=13, space_after=16)

# --------------------------------------------------------------------------
# SLIDE 10 — CLIENTS
# --------------------------------------------------------------------------
s = new_slide("Analyse clients", kicker="Comment améliorer la rétention ?", page_num=10)
add_picture_fit(s, f"{FIG}/07_pareto_clients.png", Inches(0.5), Inches(1.5), Inches(6.4), Inches(4.5))
add_picture_fit(s, f"{FIG}/08_segmentation_clients.png", Inches(7.1), Inches(1.5), Inches(5.7), Inches(4.5))
add_rect(s, Inches(0.5), Inches(6.15), Inches(12.33), Inches(1.0), GRAY_PALE)
add_text(s, Inches(0.75), Inches(6.3), Inches(11.9), Inches(0.8),
    "74 % de clients récurrents, mais 32 % des clients génèrent 80 % du CA. "
    "Le segment 'nouveaux clients' (26 %) est le principal levier de rétention.",
    size=13, color=GRAY, bold=True)

# --------------------------------------------------------------------------
# SLIDE 11 — SYNTHESE INSIGHTS
# --------------------------------------------------------------------------
s = new_slide("Synthèse des insights clés", kicker="Ce qu'il faut retenir", page_num=11)
insights = [
    "Électronique = 74,6 % du CA mais marge la plus faible et taux de retour le plus élevé",
    "Douala : plus forte croissance ET plus fort taux d'annulation — signal d'alerte opérationnel",
    "Email : meilleur ROI marketing (231x) mais largement sous-financé",
    "32 % des clients génèrent 80 % du CA ; 26 % n'ont commandé qu'une fois",
    "Profit net global estimé (~410k$) nettement inférieur au potentiel du CA brut (~2,51M$)",
]
add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), insights, size=17, space_after=22)

# --------------------------------------------------------------------------
# SLIDE 12 — 5 RECOMMANDATIONS
# --------------------------------------------------------------------------
s = new_slide("5 recommandations stratégiques", kicker="Plan d'action", page_num=12)
recos = [
    ("1", "Réduire le taux de retour Électronique", "Audit qualité produit et fournisseurs (13,8% de retours)"),
    ("2", "Corriger le problème opérationnel à Douala", "Diagnostiquer la cause du taux d'annulation (12,9%) avant d'investir"),
    ("3", "Réallouer le budget marketing", "Renforcer Email (ROI 231x), réduire Influenceur (ROI 21,6x)"),
    ("4", "Réactiver les clients à commande unique", "Campagne Email ciblée J+30 sur Mode/Beauté (26% de la base)"),
    ("5", "Prioriser Mode et Beauté pour la croissance", "Marges 45-50%, taux de retour <7,5%"),
]
y = Inches(1.65)
for num, title_r, desc in recos:
    add_rect(s, Inches(0.8), y, Inches(0.7), Inches(0.9), RED)
    add_text(s, Inches(0.8), y + Inches(0.15), Inches(0.7), Inches(0.6), num, size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y, Inches(10.5), Inches(0.45), title_r, size=16, color=GRAY, bold=True)
    add_text(s, Inches(1.7), y + Inches(0.42), Inches(10.5), Inches(0.45), desc, size=12.5, color=GRAY_LIGHT)
    y += Inches(1.02)

# --------------------------------------------------------------------------
# SLIDE 13 — CONCLUSION
# --------------------------------------------------------------------------
s = new_slide("Conclusion & prochaines étapes", kicker="Orienté action", page_num=13)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2),
    "AfriMarket dégage un CA solide (2,51 M$ sur 6 mois) mais une rentabilité bridée par sa "
    "dépendance à l'Électronique. La priorité immédiate n'est pas de générer plus de CA, mais de "
    "protéger et augmenter le profit net.",
    size=17, color=GRAY, line_spacing=1.3)
add_bullets(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(2.8),
    [
        "Court terme (0-1 mois) : diagnostiquer Douala, lancer la relance Email des clients inactifs",
        "Moyen terme (1-3 mois) : audit qualité fournisseurs Électronique, rééquilibrage budget marketing",
        "Long terme (3-6 mois) : accélérer la croissance sur Mode et Beauté",
    ], size=15, space_after=16)

# --------------------------------------------------------------------------
# SLIDE 14 — MERCI
# --------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, GRAY)
add_rect(s, 0, 0, SW, Inches(0.35), RED)
add_logo_badge(s, SW / 2, Inches(3.15), Inches(1.5))
add_text(s, Inches(1.5), Inches(4.35), Inches(10.33), Inches(0.7),
    "MERCI DE VOTRE ATTENTION", size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.0), Inches(10.33), Inches(0.5),
    "Questions & discussion", size=15, color=RGBColor(0xD9,0xD9,0xD9), align=PP_ALIGN.CENTER, italic=True)

prs.save(OUT)
print("Presentation sauvegardee ->", OUT, "|", len(prs.slides), "slides")
